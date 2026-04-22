"""Extraccion de texto de distintos formatos de documento."""
import io

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}


class UnsupportedFileTypeError(Exception):
    """El formato del fichero no esta soportado."""
    pass


def extract_text(filename: str, content: bytes) -> str:
    """
    Extrae texto plano de un fichero segun su extension.

    Args:
        filename: nombre del fichero (se usa para deducir la extension)
        content: bytes del fichero

    Returns:
        Texto extraido como string.

    Raises:
        UnsupportedFileTypeError: si el formato no esta soportado.
    """
    filename_lower = filename.lower()

    if filename_lower.endswith(".pdf"):
        return _extract_pdf(content)
    if filename_lower.endswith(".docx"):
        return _extract_docx(content)
    if filename_lower.endswith(".pptx"):
        return _extract_pptx(content)
    if filename_lower.endswith((".txt", ".md")):
        return content.decode("utf-8", errors="replace")

    raise UnsupportedFileTypeError(
        f"Formato no soportado. Permitidos: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
    )


def _extract_pdf(content: bytes) -> str:
    reader = PdfReader(io.BytesIO(content))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages).strip()


def _extract_docx(content: bytes) -> str:
    doc = DocxDocument(io.BytesIO(content))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def _extract_pptx(content: bytes) -> str:
    prs = Presentation(io.BytesIO(content))
    slides_text = []
    for slide in prs.slides:
        slide_chunks = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_chunks.append(shape.text.strip())
        if slide_chunks:
            slides_text.append("\n".join(slide_chunks))
    return "\n\n".join(slides_text).strip()